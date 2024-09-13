import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import pandas as pd
import time
import threading

CHUNK_SIZE = 5  # Adjusted chunk size
DELAY = 0.1  # Delay between processing chunks

def browse_file():
    filepath = filedialog.askopenfilename(
        filetypes=[("PowerPoint Files", "*.pptx")]
    )
    entry_file_path.delete(0, tk.END)
    entry_file_path.insert(0, filepath)

def browse_excel_file():
    filepath = filedialog.askopenfilename(
        filetypes=[("Excel Files", "*.xlsx")]
    )
    entry_excel_path.delete(0, tk.END)
    entry_excel_path.insert(0, filepath)

def load_presentation(path):
    if not path:
        messagebox.showerror("Error", "Please select a PowerPoint file.")
        return None
    return Presentation(path)

def extract_data_from_excel(file_path):
    try:
        xls = pd.ExcelFile(file_path)
        df_servers = pd.read_excel(xls, sheet_name='Servers Part of Report Cycle')
        df_format = pd.read_excel(xls, sheet_name='Format Box')

        report_names = df_servers['Report Name'].unique()
        hostname_data = df_servers[['Report Name', 'Hostname']]
        format_data = df_format[['CPU Utilization', 'Memory Utilization', 'Disk Utilization']]

        return report_names, hostname_data, format_data
    except Exception as e:
        messagebox.showerror("Error", f"An error occurred while extracting data from Excel: {e}")
        return None, None, None

def replace_draft_template(prs, new_text):
    try:
        first_slide = prs.slides[0]
        for shape in first_slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "Draft Template" in run.text:
                            run.text = run.text.replace("Draft Template", new_text)
                            run.font.color.rgb = RGBColor(0, 0, 0)
        messagebox.showinfo("Success", f"'Draft Template' replaced with '{new_text}' on the first slide.")
    except Exception as e:
        messagebox.showerror("Error", f"An error occurred: {e}")

def apply_draft_replacement(prs, draft_text):
    try:
        if draft_text:
            replace_draft_template(prs, draft_text)
    except:
        messagebox.showerror("Error", "Please enter the replacement text for 'Draft Template'.")

def apply_saw_replacements(prs, replacements):
    try:
        for slide in prs.slides:
            for shape in slide.shapes:
                process_shape(shape, replacements)
        messagebox.showinfo("Success", "Hostnames->SAW replacements applied.")
    except Exception as e:
        messagebox.showerror("Error", f"An error occurred: {e}")

def process_shape(shape, replacements):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            process_shape(s, replacements)
    elif shape.has_text_frame:
        text_frame = shape.text_frame
        replace_text_in_text_frame(text_frame, replacements)

    if shape.has_table:
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                text_frame = cell.text_frame
                replace_text_in_text_frame(text_frame, replacements)

def replace_text_in_text_frame(text_frame, replacements):
    if text_frame is not None:
        for paragraph in text_frame.paragraphs:
            full_text = ''.join([run.text for run in paragraph.runs])
            for old_text, new_text in replacements.items():
                if old_text in full_text:
                    full_text = full_text.replace(old_text, new_text)
                    for run in paragraph.runs:
                        run.text = ''
                    paragraph.runs[0].text = full_text

def set_text_color(run, rgb_color):
    run.font.color.rgb = rgb_color

def search_and_replace_value(prs, search_value, replacement_value):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        if search_value in cell.text:
                            text_frame = cell.text_frame
                            for paragraph in text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if search_value in run.text:
                                        start = run.text.find(search_value)
                                        end = start + len(search_value)
                                        run.text = run.text[:start] + replacement_value + run.text[end:]

                                        try:
                                            if float(replacement_value.strip('%')) > 85:
                                                set_text_color(run, RGBColor(255, 0, 0))
                                        except ValueError:
                                            pass
                                        return

def apply_combined_replacements(prs, combined_replacements):
    try:
        combined_replacements_list = combined_replacements.values()
        total_lines = len(combined_replacements_list)
        for start in range(0, total_lines, CHUNK_SIZE):
            chunk = list(combined_replacements_list)[start:start + CHUNK_SIZE]

            for line in chunk:
                try:
                    value_31, value_53, value_83 = line.split()
                    search_and_replace_value(prs, "31.77%", value_31)
                    search_and_replace_value(prs, "53.07%", value_53)
                    search_and_replace_value(prs, "83.07%", value_83)
                except ValueError:
                    messagebox.showerror("Error", "Each line must contain exactly three values separated by spaces.")
                    return

            progress = min(start + CHUNK_SIZE, total_lines)
            progress_label.config(text=f"Processing {progress}/{total_lines} lines...")
            root.update_idletasks()
            time.sleep(DELAY)

        messagebox.showinfo("Success", "Combined replacements applied.")
    except Exception as e:
        messagebox.showerror("Error", f"An error occurred: {e}")

def save_changes(prs, save_path):
    if prs:
        prs.save(save_path)
        messagebox.showinfo("Success", f"Changes saved to {save_path}")
    else:
        messagebox.showerror("Error", "No presentation loaded.")

def create_presentation_from_template(template_path, save_path, draft_text, replacements, combined_replacements):
    prs = load_presentation(template_path)
    if prs:
        apply_draft_replacement(prs, draft_text)
        apply_saw_replacements(prs, replacements)
        apply_combined_replacements(prs, combined_replacements)
        save_changes(prs, save_path)

def process_excel_data(excel_path, template_path):
    report_names, hostname_data, format_data = extract_data_from_excel(excel_path)
    
    if report_names is None or hostname_data is None or format_data is None:
        return
    
    for report_name in report_names:
        hostnames = hostname_data[hostname_data['Report Name'] == report_name]['Hostname'].tolist()
        format_values = format_data.iloc[0]
        cpu = format_values['CPU Utilization']
        memory = format_values['Memory Utilization']
        disk = format_values['Disk Utilization']

        draft_text = report_name
        saw_replacements = {f"SAW{i+1:02}": hostname for i, hostname in enumerate(hostnames)}
        combined_replacements = {
            "31.77%": cpu,
            "53.07%": memory,
            "83.07%": disk
        }

        save_path = f"{report_name}_modified.pptx"
        create_presentation_from_template(template_path, save_path, draft_text, saw_replacements, combined_replacements)

def start_threaded_processing():
    threading.Thread(target=process_excel_data, args=(entry_excel_path.get(), entry_file_path.get())).start()

root = tk.Tk()
root.title("PowerPoint Report Text Replacer")

notebook = ttk.Notebook(root)
notebook.grid(row=0, column=0, padx=10, pady=10)

tab1 = ttk.Frame(notebook)
notebook.add(tab1, text="Hostnames->SAW Replacements")

tk.Label(tab1, text="Select PowerPoint File:").grid(row=0, column=0, padx=10, pady=5)
entry_file_path = tk.Entry(tab1, width=50)
entry_file_path.grid(row=0, column=1, padx=10, pady=5)
tk.Button(tab1, text="Browse", command=browse_file).grid(row=0, column=2, padx=10, pady=5)

tk.Label(tab1, text="Replacement Input (one per line):").grid(row=2, column=0, padx=10, pady=5)
entry_replacements = tk.Text(tab1, width=50, height=20)
entry_replacements.grid(row=2, column=1, padx=10, pady=5)

tab_combined = ttk.Frame(notebook)
notebook.add(tab_combined, text="Combined Replacements")

tk.Label(tab_combined, text="Replacement Values for\n CPU, Memory and Disk Utilization \n (three per line, separated by spaces):").grid(row=0, column=0, padx=10, pady=5)
entry_combined = tk.Text(tab_combined, width=50, height=10)
entry_combined.grid(row=0, column=1, padx=10, pady=5)

tk.Button(tab_combined, text="Apply Combined Replacements", command=start_threaded_processing).grid(row=1, column=1, padx=10, pady=20)

tk.Label(tab1, text="Select Excel File:").grid(row=1, column=0, padx=10, pady=5)
entry_excel_path = tk.Entry(tab1, width=50)
entry_excel_path.grid(row=1, column=1, padx=10, pady=5)
tk.Button(tab1, text="Browse", command=browse_excel_file).grid(row=1, column=2, padx=10, pady=5)

progress_label = tk.Label(root, text="")
progress_label.grid(row=1, column=0, padx=10, pady=5)

root.mainloop()
